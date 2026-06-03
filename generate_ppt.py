"""Generate presentation from Conclusions*.md files."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # slide background / titles
OCEAN      = RGBColor(0x1B, 0x4F, 0x72)   # accent bars
TEAL       = RGBColor(0x17, 0xA5, 0x89)   # highlights
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
DARK_GRAY  = RGBColor(0x2C, 0x3E, 0x50)
GOLD       = RGBColor(0xF3, 0x9C, 0x12)
GREEN      = RGBColor(0x1E, 0x8B, 0x4C)
RED        = RGBColor(0xC0, 0x39, 0x2B)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=16, color=WHITE, title=None, title_color=None,
                   bullet_color=None, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            # (text, indent_level, bold, color_override)
            text, level, bold, col = item
            p.level = level
        else:
            text, level, bold, col = item, 0, False, None

        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = col if col else color

    return txBox


def section_header_slide(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    # left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), H, TEAL)

    # section number circle
    add_rect(slide, Inches(1.0), Inches(2.8), Inches(1.1), Inches(1.1), TEAL)
    add_text_box(slide, str(number), Inches(1.0), Inches(2.85), Inches(1.1), Inches(1.0),
                 font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, title, Inches(2.4), Inches(2.7), Inches(10.0), Inches(1.2),
                 font_size=38, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle, Inches(2.4), Inches(3.85), Inches(10.0), Inches(0.8),
                     font_size=20, color=TEAL, italic=True)
    return slide


def content_slide(prs, title, bg=DARK_GRAY):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, bg)
    # top accent bar
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.12), TEAL)
    # title area
    add_rect(slide, Inches(0), Inches(0.12), W, Inches(0.85), OCEAN)
    add_text_box(slide, title, Inches(0.35), Inches(0.18), Inches(12.5), Inches(0.75),
                 font_size=24, bold=True, color=WHITE)
    return slide


# ── Slide builders ──────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)

    add_rect(slide, Inches(0), Inches(0), Inches(0.3), H, TEAL)
    add_rect(slide, Inches(0), Inches(5.5), W, Inches(2.0), OCEAN)

    add_text_box(slide, "🛰️", Inches(0.8), Inches(0.4), Inches(2), Inches(1.5),
                 font_size=64, align=PP_ALIGN.LEFT, color=WHITE)

    add_text_box(slide, "Detección de Barcos",
                 Inches(0.7), Inches(1.1), Inches(12), Inches(1.1),
                 font_size=44, bold=True, color=WHITE)
    add_text_box(slide, "en Imágenes Satelitales con YOLO11",
                 Inches(0.7), Inches(2.0), Inches(12), Inches(0.8),
                 font_size=28, bold=False, color=TEAL)

    add_text_box(slide, "Leonardo Heis · Leandro Juárez · Maximiliano Miro",
                 Inches(0.7), Inches(5.65), Inches(11), Inches(0.55),
                 font_size=18, bold=True, color=WHITE)
    add_text_box(slide, "Diplomatura en Inteligencia Artificial · Computer Vision · 2026",
                 Inches(0.7), Inches(6.2), Inches(11), Inches(0.5),
                 font_size=15, color=LIGHT_GRAY)


def slide_agenda(prs):
    slide = content_slide(prs, "Estructura de la Presentación")
    items = [
        (f"  1  Motivación — Por qué este dominio",     0, False, WHITE),
        (f"  2  Dataset — Origen, composición y etiquetado", 0, False, WHITE),
        (f"  3  Detector Clásico — OpenCV",              0, False, WHITE),
        (f"  4  Entrenamiento YOLO11 — Experimentos y métricas", 0, False, WHITE),
        (f"  5  Comparación — Clásico vs. CNN vs. YOLO", 0, False, WHITE),
        (f"  6  Conclusiones — Aprendizajes y próximos pasos", 0, False, WHITE),
    ]
    y = Inches(1.2)
    for i, (text, _, bold, col) in enumerate(items):
        bg_col = OCEAN if i % 2 == 0 else DARK_GRAY
        add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(0.72), bg_col)
        add_text_box(slide, text, Inches(0.7), y + Pt(6), Inches(12.0), Inches(0.6),
                     font_size=18, bold=bold, color=WHITE)
        y += Inches(0.78)


# 1 ─ Motivación
def slides_motivacion(prs):
    section_header_slide(prs, 1, "Motivación",
                         "¿Por qué detección de barcos desde el espacio?")

    slide = content_slide(prs, "1. Motivación — El Problema")
    set_bg(slide, DARK_GRAY)

    # Left column: problem
    add_rect(slide, Inches(0.4), Inches(1.15), Inches(5.9), Inches(5.9), OCEAN)
    add_text_box(slide, "El Problema",
                 Inches(0.55), Inches(1.25), Inches(5.6), Inches(0.5),
                 font_size=20, bold=True, color=GOLD)
    bullets_l = [
        "Vigilancia marítima a escala global",
        "Control de tráfico portuario",
        "Monitoreo de pesca ilegal",
        "Respuesta ante emergencias en alta mar",
        "Cobertura continua de vastas extensiones oceánicas",
    ]
    y = Inches(1.85)
    for b in bullets_l:
        add_text_box(slide, f"▸  {b}", Inches(0.65), y, Inches(5.6), Inches(0.5),
                     font_size=15, color=WHITE)
        y += Inches(0.52)

    # Right column: challenge
    add_rect(slide, Inches(6.7), Inches(1.15), Inches(6.2), Inches(5.9), NAVY)
    add_text_box(slide, "Los Desafíos Técnicos",
                 Inches(6.85), Inches(1.25), Inches(5.9), Inches(0.5),
                 font_size=20, bold=True, color=TEAL)
    bullets_r = [
        ("Baja resolución relativa de los satélites", RED),
        ("Barcos ocupan solo 20–50 px en la imagen", RED),
        ("Variabilidad de iluminación y clima", RED),
        ("Fondos heterogéneos: puertos, espuma, olas", RED),
        ("Estructuras portuarias similares a cascos", RED),
    ]
    y = Inches(1.85)
    for b, col in bullets_r:
        add_text_box(slide, f"⚠  {b}", Inches(6.85), y, Inches(5.9), Inches(0.5),
                     font_size=15, color=col)
        y += Inches(0.52)


# 2 ─ Dataset
def slides_dataset(prs):
    section_header_slide(prs, 2, "Dataset",
                         "ShipsNet + Roboflow · 4000 chips · bounding boxes anotadas")

    slide = content_slide(prs, "2. Dataset — Origen y Composición")

    # ShipsNet card
    add_rect(slide, Inches(0.4), Inches(1.15), Inches(6.0), Inches(2.8), OCEAN)
    add_text_box(slide, "ShipsNet (Planet Labs / Kaggle)",
                 Inches(0.55), Inches(1.25), Inches(5.8), Inches(0.45),
                 font_size=17, bold=True, color=GOLD)
    for i, line in enumerate([
        "4 000 chips de 80×80 px",
        "1 000 chips con barco  |  3 000 sin barco",
        "Imágenes del satélite Planet",
        "8 escenas satelitales completas (SF Bay / Long Beach)",
    ]):
        add_text_box(slide, f"▸  {line}", Inches(0.65), Inches(1.78) + Inches(0.46 * i),
                     Inches(5.7), Inches(0.42), font_size=15, color=WHITE)

    # Roboflow card
    add_rect(slide, Inches(6.8), Inches(1.15), Inches(6.0), Inches(2.8), NAVY)
    add_text_box(slide, "Roboflow — Anotación con Bounding Boxes",
                 Inches(6.95), Inches(1.25), Inches(5.8), Inches(0.45),
                 font_size=17, bold=True, color=TEAL)
    for i, line in enumerate([
        "300 imágenes positivas subidas",
        "Auto Label (Grounding DINO / SAM)",
        "Revisión manual de anotaciones",
        "Proyecto: ship-4jnj0  |  Formato: yolov11",
    ]):
        add_text_box(slide, f"▸  {line}", Inches(6.95), Inches(1.78) + Inches(0.46 * i),
                     Inches(5.7), Inches(0.42), font_size=15, color=WHITE)

    # Splits table
    add_text_box(slide, "División Train / Valid / Test",
                 Inches(0.4), Inches(4.15), Inches(12), Inches(0.45),
                 font_size=18, bold=True, color=WHITE)

    headers = ["Split", "Imágenes", "Bounding Boxes", "Porcentaje"]
    rows = [
        ["Train",  "1 535", "7 488",  "70%"],
        ["Valid",  "512",   "2 606",  "20%"],
        ["Test",   "512",   "2 768",  "10%"],
    ]
    col_widths = [Inches(2.5), Inches(2.5), Inches(3.5), Inches(2.5)]
    col_x = [Inches(0.4), Inches(2.95), Inches(5.5), Inches(9.1)]

    for ci, (h, cw, cx) in enumerate(zip(headers, col_widths, col_x)):
        add_rect(slide, cx, Inches(4.65), cw - Inches(0.05), Inches(0.42), TEAL)
        add_text_box(slide, h, cx + Inches(0.08), Inches(4.68), cw, Inches(0.38),
                     font_size=14, bold=True, color=WHITE)

    for ri, row in enumerate(rows):
        bg = OCEAN if ri % 2 == 0 else DARK_GRAY
        for ci, (val, cw, cx) in enumerate(zip(row, col_widths, col_x)):
            add_rect(slide, cx, Inches(5.1) + Inches(0.44 * ri), cw - Inches(0.05),
                     Inches(0.42), bg)
            add_text_box(slide, val,
                         cx + Inches(0.08), Inches(5.13) + Inches(0.44 * ri),
                         cw, Inches(0.38), font_size=14, color=WHITE)

    # Etiquetado slide
    slide2 = content_slide(prs, "2. Dataset — Criterios de Etiquetado y Dificultades")
    col_items = [
        ("CRITERIOS DE ETIQUETADO", GOLD, [
            "Objeto válido: silueta visible + ≥50% del casco en el chip",
            "Barcos parcialmente nublados: incluir si silueta reconocible",
            "Estructuras portuarias fijas (diques, grúas): excluir",
            "Densidad promedio: ~5 barcos/imagen (varía 1–20+)",
        ]),
        ("DIFICULTADES ENCONTRADAS", RED, [
            "Mezcla de máscaras de segmentación → warning Ultralytics",
            "Cajas imprecisas de Auto Label → ruido de supervisión",
            "Transición clasificación → detección: anotación manual costosa",
            "Barcos muy pequeños (<10 px): difíciles de encuadrar con exactitud",
        ]),
    ]
    for ci, (title, tcol, items) in enumerate(col_items):
        x = Inches(0.4) + ci * Inches(6.5)
        add_rect(slide2, x, Inches(1.15), Inches(6.2), Inches(5.9),
                 OCEAN if ci == 0 else NAVY)
        add_text_box(slide2, title, x + Inches(0.15), Inches(1.25),
                     Inches(5.9), Inches(0.45), font_size=15, bold=True, color=tcol)
        for ii, item in enumerate(items):
            add_text_box(slide2, f"▸  {item}",
                         x + Inches(0.15), Inches(1.82) + Inches(1.05 * ii),
                         Inches(5.9), Inches(0.95),
                         font_size=14, color=WHITE)


# 3 ─ Detector Clásico
def slides_clasico(prs):
    section_header_slide(prs, 3, "Detector Clásico",
                         "OpenCV · HSV · Morfología · CNN / MobileNetV2")

    slide = content_slide(prs, "3. Detector Clásico — Pipeline OpenCV")

    steps = [
        ("1", "HSV", "Conversión al espacio de color HSV"),
        ("2", "Saturación", "Extracción del canal S (barcos coloridos)"),
        ("3", "Umbral", "Umbralización binaria adaptativa"),
        ("4", "Morfología", "Cierre (unir) + Apertura (ruido)"),
        ("5", "Contornos", "cv2.findContours"),
        ("6", "Filtrado", "Área mínima + Aspect ratio > 1.5"),
    ]
    x = Inches(0.3)
    for i, (num, label, desc) in enumerate(steps):
        add_rect(slide, x, Inches(1.35), Inches(2.05), Inches(1.5), OCEAN)
        add_text_box(slide, num, x, Inches(1.38), Inches(0.55), Inches(0.55),
                     font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x + Inches(0.05), Inches(1.9), Inches(1.95), Inches(0.4),
                     font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, x + Inches(0.05), Inches(2.35), Inches(1.95), Inches(0.45),
                     font_size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
        if i < 5:
            add_text_box(slide, "→", x + Inches(2.08), Inches(1.75), Inches(0.3), Inches(0.5),
                         font_size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        x += Inches(2.18)

    # Éxitos y fallos
    add_rect(slide, Inches(0.4), Inches(3.1), Inches(6.0), Inches(3.95), NAVY)
    add_text_box(slide, "✓  Detecta bien",
                 Inches(0.55), Inches(3.18), Inches(5.7), Inches(0.45),
                 font_size=17, bold=True, color=GREEN)
    for i, line in enumerate([
        "Barcos coloridos sobre agua oscura (alto contraste)",
        "Escenas con pocos barcos y fondo uniforme",
        "Orientaciones horizontales/verticales claras",
    ]):
        add_text_box(slide, f"▸  {line}", Inches(0.65), Inches(3.72) + Inches(0.55 * i),
                     Inches(5.6), Inches(0.5), font_size=14, color=WHITE)

    add_rect(slide, Inches(6.8), Inches(3.1), Inches(6.0), Inches(3.95), DARK_GRAY)
    add_text_box(slide, "✗  Falla en",
                 Inches(6.95), Inches(3.18), Inches(5.7), Inches(0.45),
                 font_size=17, bold=True, color=RED)
    for i, line in enumerate([
        "Barcos grises/blancos sobre mar claro",
        "Zonas portuarias con estructuras complejas",
        "Barcos muy pequeños (< 10 px de área)",
        "Sin contexto semántico: no distingue casco de dique",
    ]):
        add_text_box(slide, f"✗  {line}", Inches(6.95), Inches(3.72) + Inches(0.55 * i),
                     Inches(5.6), Inches(0.5), font_size=14, color=LIGHT_GRAY)

    # CNN slide
    slide3 = content_slide(prs, "3. CNN Baseline + MobileNetV2 — Clasificadores de Chips")

    metrics = [
        ("CNN Baseline\n(desde cero)", [
            ("Accuracy",      "0.963"),
            ("Precision",     "0.886"),
            ("Recall ★",      "0.980"),
            ("F1-Score",      "0.930"),
            ("AUC-ROC",       "0.9945"),
        ], OCEAN),
        ("MobileNetV2\n(Transfer Learning)", [
            ("Accuracy",      "0.965"),
            ("Precision ★",   "0.964"),
            ("Recall",        "0.893"),
            ("F1-Score",      "0.927"),
            ("AUC-ROC",       "0.9933"),
        ], NAVY),
    ]
    for ci, (model_name, mets, bg) in enumerate(metrics):
        x = Inches(0.4) + ci * Inches(6.5)
        add_rect(slide3, x, Inches(1.15), Inches(6.2), Inches(5.9), bg)
        add_text_box(slide3, model_name, x + Inches(0.15), Inches(1.22),
                     Inches(5.9), Inches(0.72), font_size=17, bold=True, color=GOLD)
        for mi, (metric, value) in enumerate(mets):
            row_bg = TEAL if "★" in metric else OCEAN if ci == 0 else NAVY
            add_rect(slide3, x + Inches(0.15), Inches(2.1) + Inches(0.82 * mi),
                     Inches(5.85), Inches(0.72), row_bg)
            add_text_box(slide3, metric.replace(" ★", ""),
                         x + Inches(0.25), Inches(2.14) + Inches(0.82 * mi),
                         Inches(3.5), Inches(0.65), font_size=16, color=WHITE)
            add_text_box(slide3, value,
                         x + Inches(3.9), Inches(2.14) + Inches(0.82 * mi),
                         Inches(2.0), Inches(0.65), font_size=18, bold=True,
                         color=WHITE, align=PP_ALIGN.RIGHT)

    add_text_box(slide3,
                 "Limitación: ambos modelos clasifican chips pero NO producen bounding boxes → dependen del pipeline OpenCV",
                 Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.45),
                 font_size=13, italic=True, color=GOLD)


# 4 ─ YOLO11
def slides_yolo(prs):
    section_header_slide(prs, 4, "Entrenamiento YOLO11",
                         "yolo11m.pt · 3 experimentos · mAP@50 hasta 0.616")

    # Experimentos
    slide = content_slide(prs, "4. YOLO11 — Experimentos de Entrenamiento")

    headers = ["Corrida", "Epochs", "imgsz", "cls", "Tiempo", "mAP@50", "Precisión", "Recall"]
    rows = [
        ["v1  baseline",  "30",  "640", "0.5", "0.62 h", "0.463", "0.507", "0.401"],
        ["v2  +epochs",   "100", "640", "1.5", "1.82 h", "0.616", "0.748", "0.468"],
        ["v3  +imgsz",    "120", "960", "0.5", "3.91 h", "0.611", "0.748", "0.469"],
    ]
    col_widths = [Inches(2.2), Inches(0.9), Inches(0.9), Inches(0.8),
                  Inches(1.0), Inches(1.1), Inches(1.2), Inches(1.1)]
    col_x = [sum(col_widths[:i]) + Inches(0.4) for i in range(len(col_widths))]

    for ci, (h, cw, cx) in enumerate(zip(headers, col_widths, col_x)):
        add_rect(slide, cx, Inches(1.2), cw - Inches(0.03), Inches(0.45), TEAL)
        add_text_box(slide, h, cx + Inches(0.05), Inches(1.23),
                     cw, Inches(0.4), font_size=12, bold=True, color=WHITE)

    row_colors = [OCEAN, NAVY, OCEAN]
    for ri, (row, bg) in enumerate(zip(rows, row_colors)):
        for ci, (val, cw, cx) in enumerate(zip(row, col_widths, col_x)):
            cell_bg = TEAL if (ri == 1 and ci in (5, 6)) else bg
            add_rect(slide, cx, Inches(1.68) + Inches(0.52 * ri),
                     cw - Inches(0.03), Inches(0.5), cell_bg)
            add_text_box(slide, val,
                         cx + Inches(0.05), Inches(1.71) + Inches(0.52 * ri),
                         cw, Inches(0.44), font_size=12,
                         bold=(ri == 1 and ci in (5, 6)), color=WHITE)

    # Análisis
    analyses = [
        ("v1 → v2: mayor impacto",
         "Duplicar epochs (30→100) + aumentar cls (0.5→1.5): mAP@50 sube 33% (0.463→0.616). "
         "v1 estaba claramente subajustado — la pérdida de validación seguía descendiendo al finalizar.",
         GREEN),
        ("v2 → v3: ganancia marginal",
         "Mayor resolución (640→960 px) mejora barcos pequeños, pero batch cae a 3 imágenes "
         "(AutoBatch al 52% VRAM). mAP pasa de 0.616 a 0.611 — no hay mejora neta.",
         GOLD),
    ]
    for ai, (title, body, col) in enumerate(analyses):
        add_rect(slide, Inches(0.4), Inches(3.5) + Inches(1.7 * ai),
                 Inches(12.5), Inches(1.55), NAVY if ai == 0 else DARK_GRAY)
        add_text_box(slide, title,
                     Inches(0.55), Inches(3.57) + Inches(1.7 * ai),
                     Inches(12.0), Inches(0.45), font_size=15, bold=True, color=col)
        add_text_box(slide, body,
                     Inches(0.55), Inches(4.05) + Inches(1.7 * ai),
                     Inches(12.0), Inches(0.85), font_size=13, color=WHITE)

    # Métricas mejor modelo
    slide2 = content_slide(prs, "4. YOLO11 — Métricas del Mejor Modelo (v3, test set)")

    metrics_vals = [
        ("mAP@50",          "0.611", TEAL),
        ("mAP@50-95",       "0.333", OCEAN),
        ("Precisión",       "0.748", OCEAN),
        ("Recall",          "0.469", RED),
        ("Inferencia GPU",  "25.8 ms/img", OCEAN),
    ]
    for mi, (metric, value, bg) in enumerate(metrics_vals):
        x = Inches(0.4) + mi * Inches(2.55)
        add_rect(slide2, x, Inches(1.3), Inches(2.4), Inches(2.4), bg)
        add_text_box(slide2, metric, x + Inches(0.1), Inches(1.4),
                     Inches(2.2), Inches(0.55), font_size=15, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide2, value, x + Inches(0.1), Inches(2.0),
                     Inches(2.2), Inches(0.8), font_size=28, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)

    insights = [
        ("mAP@50 (0.611) vs. mAP@50-95 (0.333)",
         "El modelo localiza bien los barcos (IoU > 0.5), "
         "pero no ajusta perfectamente los bordes de la caja (falla a IoU > 0.75). "
         "Las anotaciones automáticas de Auto Label introducen ruido en las cajas.",
         OCEAN),
        ("Recall bajo (0.47) — el problema prioritario",
         "Para vigilancia marítima, perder un barco (falso negativo) es peor que una falsa alarma. "
         "Se recomienda operar con conf=0.20–0.30 en lugar del default 0.50, "
         "sacrificando precisión para recuperar recall.",
         NAVY),
        ("Matriz de confusión",
         "Detecta barcos en alta densidad portuaria. "
         "Confunde barcos pequeños con fondo (↑ FN). "
         "Pocas detecciones espurias sobre agua abierta (Precisión 0.75).",
         DARK_GRAY),
    ]
    for ii, (title, body, bg) in enumerate(insights):
        add_rect(slide2, Inches(0.4), Inches(3.95) + Inches(1.12 * ii),
                 Inches(12.5), Inches(1.05), bg)
        add_text_box(slide2, title,
                     Inches(0.55), Inches(4.0) + Inches(1.12 * ii),
                     Inches(12.0), Inches(0.38), font_size=14, bold=True, color=TEAL)
        add_text_box(slide2, body,
                     Inches(0.55), Inches(4.4) + Inches(1.12 * ii),
                     Inches(12.0), Inches(0.6), font_size=13, color=WHITE)


# 5 ─ Comparación
def slides_comparacion(prs):
    section_header_slide(prs, 5, "Comparación de Enfoques",
                         "Detector Clásico · CNN/MobileNetV2 · YOLO11")

    slide = content_slide(prs, "5. Comparación — Tabla Resumen")

    dims = [
        "Salida del modelo",
        "Precisión cualitativa",
        "Requisitos computacionales",
        "Tiempo de desarrollo",
        "Generalización",
    ]
    clasico_vals = [
        "Contornos aproximados",
        "Baja: falla grises, puertos, objetos pequeños",
        "Mínimos: CPU en tiempo real",
        "Rápido (~2 hs); difícil de mejorar",
        "Muy baja: parámetros calibrados por escena",
    ]
    cnn_vals = [
        "Clase + probabilidad (sin localización)",
        "Alta clasificación chips: acc 96%, AUC 0.99",
        "Bajos: inferencia CPU (9.5 MB)",
        "Moderado: train en minutos sobre chips",
        "Media: no generaliza a otras resoluciones",
    ]
    yolo_vals = [
        "Bounding boxes + coordenadas + confianza",
        "Media-alta: mAP@50=0.61, detección end-to-end",
        "GPU para entrenar; 25.8 ms/img inferencia",
        "Significativo: v1=0.6h, v2=1.8h, v3=3.9h",
        "Media: necesita más datos para ↑ recall",
    ]

    headers2 = ["Dimensión", "Clásico (OpenCV)", "CNN / MobileNetV2", "YOLO11m"]
    col_widths2 = [Inches(2.8), Inches(3.1), Inches(3.3), Inches(3.3)]
    header_colors = [DARK_GRAY, OCEAN, OCEAN, TEAL]
    col_x2 = [sum(col_widths2[:i]) + Inches(0.2) for i in range(len(col_widths2))]

    for ci, (h, cw, cx, hc) in enumerate(zip(headers2, col_widths2, col_x2, header_colors)):
        add_rect(slide, cx, Inches(1.2), cw - Inches(0.05), Inches(0.45), hc)
        add_text_box(slide, h, cx + Inches(0.07), Inches(1.23),
                     cw, Inches(0.4), font_size=13, bold=True, color=WHITE)

    for ri, (dim, cv, cnv, yv) in enumerate(zip(dims, clasico_vals, cnn_vals, yolo_vals)):
        row_bg = DARK_GRAY if ri % 2 == 0 else NAVY
        row_h = Inches(1.05)
        for ci, (val, cw, cx) in enumerate(zip([dim, cv, cnv, yv], col_widths2, col_x2)):
            cell_bg = OCEAN if (ci == 3) else row_bg
            add_rect(slide, cx, Inches(1.68) + row_h * ri, cw - Inches(0.05), row_h, cell_bg)
            add_text_box(slide, val, cx + Inches(0.07), Inches(1.72) + row_h * ri,
                         cw - Inches(0.1), row_h - Inches(0.1),
                         font_size=11, color=WHITE, bold=(ci == 0))

    # Conclusión comparativa
    slide2 = content_slide(prs, "5. Comparación — Análisis Cualitativo")
    set_bg(slide2, DARK_GRAY)

    conclusions_comp = [
        ("Clásico como baseline rápido", OCEAN,
         "Implementación en ~2 hs sin datos etiquetados. Útil para validar que el problema "
         "tiene señal visual. Techo bajo: los parámetros calibrados fallan en otra condición "
         "de iluminación o escena sin re-calibración."),
        ("CNN/MobileNetV2: excelente clasificación, sin localización", NAVY,
         "Accuracy 96%, AUC-ROC 0.99. Ideal cuando solo se necesita saber si hay un barco en el chip. "
         "Limitación crítica: no produce coordenadas. Depende del pipeline clásico para extraer candidatos."),
        ("YOLO11: el único end-to-end con localización", TEAL,
         "mAP@50=0.61. Detecta barcos con coordenadas precisas en imágenes completas. "
         "Limitación actual: recall 0.47 por falta de datos. "
         "Con 3000+ imágenes de train el recall subiría significativamente."),
    ]
    for ci, (title, bg, body) in enumerate(conclusions_comp):
        add_rect(slide2, Inches(0.4), Inches(1.2) + Inches(2.0 * ci),
                 Inches(12.5), Inches(1.85), bg)
        add_text_box(slide2, title,
                     Inches(0.55), Inches(1.28) + Inches(2.0 * ci),
                     Inches(12.0), Inches(0.5), font_size=17, bold=True, color=GOLD)
        add_text_box(slide2, body,
                     Inches(0.55), Inches(1.82) + Inches(2.0 * ci),
                     Inches(12.0), Inches(1.0), font_size=14, color=WHITE)


# 6 ─ Conclusiones
def slides_conclusiones(prs):
    section_header_slide(prs, 6, "Conclusiones",
                         "¿Qué aprendimos? ¿Qué haríamos diferente?")

    slide = content_slide(prs, "6. Conclusiones — Lo que Aprendimos")

    learnings = [
        ("1", "El etiquetado es el cuello de botella",
         "Auto Label + revisión manual tomó más tiempo que el entrenamiento. "
         "Cajas imprecisas se propagan directamente al recall.", GOLD),
        ("2", "Los epochs importan más que la arquitectura",
         "v1→v2: mAP sube 33% solo por entrenar más epochs (30→100). "
         "El modelo v1 estaba claramente subajustado.", GREEN),
        ("3", "La resolución ayuda, pero el batch la limita",
         "640→960 px mejoró la representación de objetos pequeños, "
         "pero batch cayó a 3 en una GPU de 8 GB. Ganancia marginal, costo alto.", TEAL),
        ("4", "El recall debe guiar todas las decisiones",
         "No detectar un barco es peor que una falsa alarma. "
         "Operar con conf=0.20–0.30, no 0.50.", RED),
        ("5", "El clásico falla por falta de contexto semántico",
         "Umbralización y morfología no pueden aprender qué es un barco. "
         "Solo capturan propiedades espectrales locales.", OCEAN),
    ]
    for li, (num, title, body, col) in enumerate(learnings):
        x = Inches(0.3)
        y = Inches(1.18) + Inches(1.2 * li)
        add_rect(slide, x, y, Inches(0.55), Inches(1.08), col)
        add_text_box(slide, num, x, y + Inches(0.25),
                     Inches(0.55), Inches(0.55), font_size=22, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, x + Inches(0.58), y, Inches(12.1), Inches(1.08), NAVY)
        add_text_box(slide, title,
                     x + Inches(0.72), y + Inches(0.06),
                     Inches(11.8), Inches(0.42), font_size=15, bold=True, color=col)
        add_text_box(slide, body,
                     x + Inches(0.72), y + Inches(0.5),
                     Inches(11.8), Inches(0.52), font_size=13, color=WHITE)

    # Qué haríamos diferente
    slide2 = content_slide(prs, "6. Conclusiones — ¿Qué Haríamos Diferente?")

    next_steps = [
        ("Anotar más imágenes", "IMPACTO ALTO",
         "~1500 imágenes de train es un techo. Duplicar (≥3000) es la mejora de mayor impacto en recall.",
         TEAL, GREEN),
        ("Revisar anotaciones de Auto Label", "IMPACTO ALTO",
         "Cajas imprecisas limitan la convergencia. Revisión manual de calidad mejora recall sin más datos.",
         TEAL, GREEN),
        ("Integrar SAHI desde el inicio", "IMPACTO ALTO",
         "Slicing Aided Hyper Inference para escenas de alta resolución. "
         "Transforma 0 detecciones en SF Bay a detecciones reales.",
         OCEAN, GOLD),
        ("Agregar negativos difíciles de SF Bay", "IMPACTO MEDIO",
         "Chips de espuma, diques y estructuras portuarias como ejemplos negativos "
         "para reducir falsos positivos en escenas reales.",
         OCEAN, GOLD),
        ("Explorar imgsz=1280 con batch=1", "IMPACTO MEDIO",
         "Evaluar si el trade-off resolución/batch favorece barcos muy pequeños.",
         NAVY, LIGHT_GRAY),
    ]
    for ni, (title, impact, body, bg, tcol) in enumerate(next_steps):
        row_top = Inches(1.2) + Inches(1.18 * ni)
        add_rect(slide2, Inches(0.3), row_top, Inches(12.7), Inches(1.1), bg)
        add_text_box(slide2, impact,
                     Inches(9.8), row_top + Inches(0.06),
                     Inches(3.0), Inches(0.38), font_size=11, bold=True,
                     color=tcol, align=PP_ALIGN.RIGHT)
        add_text_box(slide2, title,
                     Inches(0.45), row_top + Inches(0.06),
                     Inches(9.2), Inches(0.4), font_size=15, bold=True, color=WHITE)
        add_text_box(slide2, body,
                     Inches(0.45), row_top + Inches(0.5),
                     Inches(12.1), Inches(0.55), font_size=13, color=LIGHT_GRAY)

    # Experiencia Leandro / Colab
    slide3 = content_slide(prs, "6. Experiencia Personal — Ejecución en Google Colab (LJ)")

    add_rect(slide3, Inches(0.3), Inches(1.15), Inches(12.7), Inches(0.65), OCEAN)
    add_text_box(slide3,
                 "Todas las ejecuciones realizadas en Google Colab (GPU T4) — sin GPU local en macOS",
                 Inches(0.45), Inches(1.22), Inches(12.4), Inches(0.5),
                 font_size=16, bold=True, color=WHITE)

    versions = [
        ("Primera versión", RED,
         "Inferencia directa sobre escenas SF Bay (~2500×1700 px)\n"
         "→  RESULTADO: 0 detecciones\n"
         "→  CAUSA: barcos de 10–15 px invisibles para modelo entrenado con chips de 80 px"),
        ("Segunda versión + SAHI", GREEN,
         "Slicing en tiles de 320×320 px con 20% overlap, fusión con NMS\n"
         "→  RESULTADO: barcos detectados en SF Bay\n"
         "→  PENDIENTE: falsos positivos en espuma y estructuras portuarias"),
    ]
    for vi, (title, col, body) in enumerate(versions):
        x = Inches(0.3) + vi * Inches(6.5)
        add_rect(slide3, x, Inches(2.0), Inches(6.2), Inches(3.5), NAVY)
        add_rect(slide3, x, Inches(2.0), Inches(6.2), Inches(0.5), col)
        add_text_box(slide3, title, x + Inches(0.1), Inches(2.02),
                     Inches(5.9), Inches(0.44), font_size=16, bold=True, color=WHITE)
        add_text_box(slide3, body, x + Inches(0.12), Inches(2.6),
                     Inches(5.85), Inches(2.8), font_size=14, color=WHITE)

    add_rect(slide3, Inches(0.3), Inches(5.7), Inches(12.7), Inches(1.45), DARK_GRAY)
    add_text_box(slide3, "Notebook macOS (ship_detection_macos.ipynb)",
                 Inches(0.45), Inches(5.78), Inches(12.0), Inches(0.42),
                 font_size=14, bold=True, color=TEAL)
    add_text_box(slide3,
                 "Adaptada para Apple Silicon (MPS) o CPU · credenciales vía .env · "
                 "batch=8 (MPS)/4 (CPU) · orientada a inferencia local con pesos entrenados en Colab",
                 Inches(0.45), Inches(6.22), Inches(12.0), Inches(0.85),
                 font_size=13, color=WHITE)


def slide_cierre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.3), H, TEAL)
    add_rect(slide, Inches(0), Inches(5.8), W, Inches(1.7), OCEAN)

    add_text_box(slide, "Gracias",
                 Inches(1.0), Inches(1.8), Inches(11), Inches(1.3),
                 font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "¿Preguntas?",
                 Inches(1.0), Inches(3.1), Inches(11), Inches(0.8),
                 font_size=30, color=TEAL, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "Leonardo Heis · Leandro Juárez · Maximiliano Miro",
                 Inches(1.0), Inches(5.95), Inches(11), Inches(0.5),
                 font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "Diplomatura en IA · Computer Vision · 2026",
                 Inches(1.0), Inches(6.45), Inches(11), Inches(0.45),
                 font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_cover(prs)
    slide_agenda(prs)
    slides_motivacion(prs)
    slides_dataset(prs)
    slides_clasico(prs)
    slides_yolo(prs)
    slides_comparacion(prs)
    slides_conclusiones(prs)
    slide_cierre(prs)

    out = "presentacion_ship_detection.pptx"
    prs.save(out)
    print(f"✅  Presentación guardada en: {out}")
    print(f"    Slides generados: {len(prs.slides)}")


if __name__ == "__main__":
    main()
